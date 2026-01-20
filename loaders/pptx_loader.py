from pathlib import Path
from pptx import Presentation

def load_pptx_file(file_path: str) -> list:
    """
    Extract text from PowerPoint (.pptx) files.
    Returns list of slides with text content and metadata.
    """
    path = Path(file_path)

    if not path.exists():
        raise FileNotFoundError(f"{file_path} not found")

    try:
        prs = Presentation(path)
    except Exception as e:
        raise ValueError(f"Error reading PPTX file {file_path}: {e}")

    docs = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_content = []

        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text)

            # Extract text from tables
            if shape.has_table:
                table = shape.table
                table_text = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_text.append(" | ".join(row_data))
                if table_text:
                    slide_content.append("TABLE:\n" + "\n".join(table_text))

        # Only add non-empty slides
        if slide_content:
            docs.append({
                "content": "\n\n".join(slide_content),
                "metadata": {
                    "source": path.name,
                    "slide": slide_num,
                    "type": "pptx",
                    "total_slides": len(prs.slides)
                }
            })

    return docs
